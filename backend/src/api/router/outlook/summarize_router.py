from fastapi import APIRouter, Depends, HTTPException, status, Request
from typing import List, Optional, Any, Dict, Union
from pydantic import BaseModel, Field
from enum import Enum
import os
import sys
import json
import base64
import io
from typing import Optional, Dict, Any
from fastapi import Request, Depends, HTTPException, status
from PyPDF2 import PdfReader
from docx import Document  # python-docx
from openpyxl import load_workbook  # Excel
from pptx import Presentation  # PowerPoint
from .email_config import SupportedLanguage

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', '..')))
from src.services.auth.middleware.auth_firebase import get_current_user
from src.services.llm.llm import LLM
from src.core.logger import log
from src.api.utils.config_loader import get_summarize_config, is_style_analysis_enabled

logger = log.bind(name="src.api.outlook.summarize_router")

# Models
class SummaryType(str, Enum):
    """Types of summary available"""
    CONCISE = "concise"
    DETAILED = "detailed"
    BULLET_POINTS = "bullet_points"
    ACTION_ITEMS = "action_items"

class Source(BaseModel):
    """Source document used in RAG"""
    filename: str
    page_number: Optional[int] = None
    content: str

class ProcessingMode(str, Enum):
    """Processing mode used for the request"""
    CLIENT = "client"
    SERVER = "server"
    HYBRID = "hybrid"
    
class SummarizeRequest(BaseModel):
    """Request model for summarization"""
    authToken: Optional[str] = None
    userId: str
    
    # Email fields
    subject: Optional[str] = None
    from_email: Optional[str] = Field(None, alias="from")
    body: Optional[str] = None  # Now optional to support file-only requests
    
    # Configuration
    language: SupportedLanguage = SupportedLanguage.ENGLISH.name
    summary_type: SummaryType = SummaryType.CONCISE
    use_rag: Optional[bool] = False
    processing_mode: Optional[ProcessingMode] = ProcessingMode.SERVER
    
    # File summarization
    file_name: Optional[str] = None
    file_content: Optional[str] = None  # Base64 encoded file content (server-side processing)
    file_type: Optional[str] = None
    extracted_text: Optional[str] = None  # Pre-extracted text from client (client-side processing)

class SummarizeResponse(BaseModel):
    """Response model for summarization"""
    summary: str
    sources: List[Union[Source, str]] = []
    summary_type: SummaryType
    temperature: Optional[float] = None
    model: Optional[str] = None
    use_retrieval: Optional[bool] = None

# Router
router = APIRouter(prefix="/outlook", tags=["outlook"])

# Prompt templates for different summary types
SUMMARY_SYSTEM_PROMPTS = {
    SummaryType.CONCISE: """You are an AI assistant that summarizes content concisely. 
{language_instruction}

Content to summarize:
{content}

Create a concise summary (about 2-3 sentences) of the content that captures the most important information. Focus on key points only.
The summary should be:
- Brief and to the point
- Written in {language_name}
- Include only essential information
""",
    
    SummaryType.DETAILED: """You are an AI assistant that creates comprehensive summaries. 
{language_instruction}

Content to summarize:
{content}

Create a detailed summary of the content that captures all important information, context, and nuances. This should be a comprehensive overview.
The summary should be:
- Thorough and detailed (around 5-8 sentences)
- Written in {language_name}
- Include all important information
- Well-structured with logical flow
""",
    
    SummaryType.BULLET_POINTS: """You are an AI assistant that summarizes content in bullet point format. 
{language_instruction}

Content to summarize:
{content}

Create a bullet-point summary of the content that organizes the key information in an easy-to-read list format.
The summary should be:
- Written in {language_name}
- Formatted as bullet points (•)
- Include all important information
- Organized by topic or importance
- Brief but comprehensive
""",
    
    SummaryType.ACTION_ITEMS: """You are an AI assistant that identifies action items from content. 
{language_instruction}

Content to summarize:
{content}

Extract all action items, tasks, requests, and commitments from the content. Format them as a clear action item list.
The summary should be:
- Written in {language_name}
- Focused only on action items and tasks
- Formatted as a numbered list (1., 2., etc.)
- Specify who needs to do what (if mentioned)
- Include any relevant deadlines mentioned
"""
}
# Create API endpoints
router = APIRouter(prefix="/outlook", tags=["outlook"])

@router.get("/summarize/health")
async def health_check():
    """Health check endpoint for Outlook summarization router"""
    return {"status": "healthy", "service": "outlook-summarization"}

@router.post("/summarize", response_model=SummarizeResponse)
async def summarize(
    request_data: SummarizeRequest,
    request: Request,
    current_user: Optional[Dict[str, Any]] = Depends(get_current_user)
):
    """Summarize email content or file attachment using direct LLM"""
    try:
        # Get configuration settings
        config = get_summarize_config()
        
        # Log the request (excluding sensitive content)
        logger.info(f"Received summarize request for user {request_data.userId}, language: {request_data.language}, summary type: {request_data.summary_type}, processing mode: {request_data.processing_mode}")
        
        # Verify user
        if current_user is None:
            # For development and testing, allow processing without auth
            logger.warning(f"No authenticated user found, proceeding with userId from request: {request_data.userId}")
            current_user = {"uid": request_data.userId}
        elif current_user.get("uid") != request_data.userId:
            logger.warning(f"User ID mismatch: {current_user.get('uid')} vs {request_data.userId}")
        
        # Determine what content we're working with
        is_file_request = request_data.file_name is not None
        is_email_request = request_data.body is not None
        
        content_to_summarize = ""
        
        # Get the content to summarize based on processing mode
        if is_file_request:
            logger.info(f"Processing file server-side for {request_data.file_name}")
            try:
                # Decode base64 into raw bytes
                file_bytes = base64.b64decode(request_data.file_content)
                file_name = request_data.file_name.lower()

                if file_name.endswith(".pdf"):
                    # Handle PDF
                    pdf_file = io.BytesIO(file_bytes)
                    reader = PdfReader(pdf_file)
                    content_to_summarize = ""
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            content_to_summarize += text + "\n"

                elif file_name.endswith(".docx"):
                    # Handle Word documents
                    doc_file = io.BytesIO(file_bytes)
                    doc = Document(doc_file)
                    content_to_summarize = "\n".join(
                        [para.text for para in doc.paragraphs if para.text.strip()]
                    )

                elif file_name.endswith(".xlsx"):
                    # Handle Excel files
                    excel_file = io.BytesIO(file_bytes)
                    wb = load_workbook(excel_file, data_only=True)
                    content_to_summarize = ""
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        content_to_summarize += f"\n--- Sheet: {sheet} ---\n"
                        for row in ws.iter_rows(values_only=True):
                            row_text = " | ".join([str(cell) for cell in row if cell is not None])
                            if row_text.strip():
                                content_to_summarize += row_text + "\n"

                elif file_name.endswith(".pptx"):
                    # Handle PowerPoint files
                    ppt_file = io.BytesIO(file_bytes)
                    prs = Presentation(ppt_file)
                    content_to_summarize = ""
                    for i, slide in enumerate(prs.slides, start=1):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                if shape.text.strip():
                                    slide_text.append(shape.text.strip())
                        if slide_text:
                            content_to_summarize += f"\n--- Slide {i} ---\n" + "\n".join(slide_text)

                else:
                    # Assume plain-text-like file
                    content_to_summarize = file_bytes.decode("utf-8", errors="replace")

                logger.info("File content successfully extracted")
                if content_to_summarize is None or content_to_summarize == "":
                    return SummarizeResponse(
                        summary="Impossible to summarize this file",
                        sources=[],  # No sources when using direct LLM
                        summary_type=request_data.summary_type,
                        temperature=0.3,
                        model="",  # Use the model name from LLM client
                        use_retrieval=False  # Direct LLM doesn't use retrieval
                    )

            except Exception as e:
                logger.error(f"Error processing file {file_name}: {str(e)}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Failed to extract text from {file_name}: {str(e)}"
                )
        elif is_email_request:
            # Email summarization
            logger.info("Processing email content")
            content_to_summarize = request_data.body
        else:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Request must include either email body or file information"
            )
            
        # Build the prompt based on the request type
        if is_file_request:
            # File summarization
            system_prompt = f"You are a helpful assistant that summarizes document content. Provide a {request_data.summary_type.value} summary of the document content."
            
            user_prompt = f"""Document: {request_data.file_name} ({request_data.file_type})

Content:
{content_to_summarize}

Provide a {request_data.summary_type.value} summary of this document in this language: {request_data.language}."""
        else:
            # Email summarization
            system_prompt = f"You are a helpful assistant that summarizes email content. Provide a {request_data.summary_type.value} summary of the email content."
            
            # Create a prompt with the email content
            from_part = f"From: {request_data.from_email}\n" if request_data.from_email else ""
            subject_part = f"Subject: {request_data.subject}\n" if request_data.subject else ""
            
            user_prompt = f"""{from_part}{subject_part}
Email Content:
{content_to_summarize}

Provide a {request_data.summary_type.value} summary of this email in this language: {request_data.language}."""
        
        # Add summary type specific instructions
        if request_data.summary_type == SummaryType.CONCISE:
            user_prompt += "\n\nKeep the summary brief and to the point, focusing only on the most important information."
        elif request_data.summary_type == SummaryType.DETAILED:
            user_prompt += "\n\nProvide a comprehensive summary including all key points, details, and context."
        elif request_data.summary_type == SummaryType.BULLET_POINTS:
            user_prompt += "\n\nFormat the summary as bullet points, with each point representing a distinct piece of information."
        elif request_data.summary_type == SummaryType.ACTION_ITEMS:
            user_prompt += "\n\nExtract and list all action items, tasks, requests, and deadlines mentioned in the content."
        
        # Get user's style analysis for personalization
        user_id = current_user.get("uid") if current_user else "anonymous"
        
        # Create message for LLM service
        messages = [{"role": "user", "content": content_to_summarize}]
        
        # Initialize LLM client with configuration temperature
        llm_client = LLM(temperature=config['default_temperature'])  # Use config temperature for summaries
        
        # Call LLM service to generate summary
        try:
            # Use the chat method with system prompt and user message
            summary = await llm_client.chat(
                messages=messages,
                system_prompt=system_prompt
            )
            
            # For logging purposes
            logger.info(f"Generated summary using LLM service for user {current_user['uid']}")
        except Exception as llm_error:
            logger.error(f"Error calling LLM service: {str(llm_error)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Error generating summary: {str(llm_error)}"
            )
        
        # Return the summarized content
        return SummarizeResponse(
            summary=summary,
            sources=[],  # No sources when using direct LLM
            summary_type=request_data.summary_type,
            temperature=config['default_temperature'],
            model=llm_client.model or config['default_model'],  # Use the model name from LLM client or config
            use_retrieval=False  # Direct LLM doesn't use retrieval
        )
        
    except Exception as e:
        logger.error(f"Error generating summary for user {request_data.userId}: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate summary: {str(e)}"
        )


