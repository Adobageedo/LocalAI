"""
File Processors for Attachments
Handles different file types: Office documents, images, PDFs, text files
"""

import os
import base64
import mimetypes
from typing import Optional, Dict, Any, List
from pathlib import Path
import tempfile

# Office document processing
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# PDF processing
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Image processing
from PIL import Image
import io


class FileProcessor:
    """Base class for file processors"""
    
    SUPPORTED_TEXT_FORMATS = ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.css', '.js', '.ts', '.log']
    SUPPORTED_IMAGE_FORMATS = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']
    SUPPORTED_OFFICE_FORMATS = ['.docx', '.xlsx', '.pptx']
    SUPPORTED_PDF_FORMATS = ['.pdf']
    
    @staticmethod
    def get_file_type(filename: str) -> str:
        """Determine file type from extension"""
        ext = Path(filename).suffix.lower()
        
        if ext in FileProcessor.SUPPORTED_TEXT_FORMATS:
            return 'text'
        elif ext in FileProcessor.SUPPORTED_IMAGE_FORMATS:
            return 'image'
        elif ext in FileProcessor.SUPPORTED_OFFICE_FORMATS:
            return 'office'
        elif ext in FileProcessor.SUPPORTED_PDF_FORMATS:
            return 'pdf'
        else:
            return 'unknown'
    
    @staticmethod
    def get_mime_type(filename: str) -> str:
        """Get MIME type for file"""
        mime_type, _ = mimetypes.guess_type(filename)
        return mime_type or 'application/octet-stream'


class TextFileProcessor:
    """Process text-based files"""
    
    @staticmethod
    def extract_text(file_content: bytes, filename: str) -> str:
        """Extract text from text files"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    return file_content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # Fallback: ignore errors
            return file_content.decode('utf-8', errors='ignore')
        except Exception as e:
            raise ValueError(f"Failed to extract text from {filename}: {str(e)}")


class ImageProcessor:
    """Process image files for GPT Vision"""
    
    @staticmethod
    def prepare_for_vision(file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Prepare image for GPT Vision API
        Returns dict with base64 encoded image and metadata
        """
        try:
            # Validate it's a valid image
            img = Image.open(io.BytesIO(file_content))
            
            # Get image info
            width, height = img.size
            format_name = img.format or 'UNKNOWN'
            
            # Convert to RGB if necessary (for PNG with transparency, etc.)
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            
            # Resize if too large (GPT-4 Vision has limits)
            max_size = 2048
            if width > max_size or height > max_size:
                ratio = min(max_size / width, max_size / height)
                new_size = (int(width * ratio), int(height * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            
            # Convert to bytes
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=85)
            buffer.seek(0)
            image_bytes = buffer.read()
            
            # Encode to base64
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            
            # Determine MIME type
            mime_type = FileProcessor.get_mime_type(filename)
            if mime_type.startswith('image/'):
                data_url = f"data:{mime_type};base64,{base64_image}"
            else:
                data_url = f"data:image/jpeg;base64,{base64_image}"
            
            return {
                'type': 'image',
                'filename': filename,
                'data_url': data_url,
                'width': img.size[0],
                'height': img.size[1],
                'format': format_name,
                'size_bytes': len(image_bytes)
            }
        
        except Exception as e:
            raise ValueError(f"Failed to process image {filename}: {str(e)}")
    
    @staticmethod
    def create_vision_message_content(image_data: Dict[str, Any], user_prompt: str = None) -> List[Dict[str, Any]]:
        """
        Create message content for GPT Vision API
        Returns list with text and image_url parts
        """
        content = []
        
        # Add text prompt if provided
        if user_prompt:
            content.append({
                "type": "text",
                "text": user_prompt
            })
        
        # Add image
        content.append({
            "type": "image_url",
            "image_url": {
                "url": image_data['data_url'],
                "detail": "high"  # or "low" for faster processing
            }
        })
        
        return content


class OfficeFileProcessor:
    """Process Microsoft Office files"""
    
    @staticmethod
    def extract_from_docx(file_content: bytes, filename: str) -> str:
        """Extract text from DOCX files"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is not installed. Run: pip install python-docx")
        
        try:
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            
            try:
                doc = DocxDocument(tmp_path)
                
                # Extract text from paragraphs
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                
                # Extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            text_parts.append(row_text)
                
                return '\n\n'.join(text_parts)
            
            finally:
                os.unlink(tmp_path)
        
        except Exception as e:
            raise ValueError(f"Failed to extract text from DOCX {filename}: {str(e)}")
    
    @staticmethod
    def extract_from_xlsx(file_content: bytes, filename: str) -> str:
        """Extract data from XLSX files"""
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl is not installed. Run: pip install openpyxl")
        
        try:
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            
            try:
                workbook = openpyxl.load_workbook(tmp_path, data_only=True)
                
                text_parts = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_parts.append(f"## Sheet: {sheet_name}\n")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                        if row_text.strip():
                            text_parts.append(row_text)
                    
                    text_parts.append("")  # Blank line between sheets
                
                return '\n'.join(text_parts)
            
            finally:
                os.unlink(tmp_path)
        
        except Exception as e:
            raise ValueError(f"Failed to extract data from XLSX {filename}: {str(e)}")
    
    @staticmethod
    def extract_from_pptx(file_content: bytes, filename: str) -> str:
        """Extract text from PPTX files"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. Run: pip install python-pptx")
        
        try:
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            
            try:
                prs = Presentation(tmp_path)
                
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, start=1):
                    text_parts.append(f"## Slide {slide_num}\n")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                    
                    text_parts.append("")  # Blank line between slides
                
                return '\n'.join(text_parts)
            
            finally:
                os.unlink(tmp_path)
        
        except Exception as e:
            raise ValueError(f"Failed to extract text from PPTX {filename}: {str(e)}")


class PDFProcessor:
    """Process PDF files"""
    
    @staticmethod
    def extract_text(file_content: bytes, filename: str) -> str:
        """Extract text from PDF files"""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 is not installed. Run: pip install PyPDF2")
        
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"## Page {page_num}\n{text}")
            
            return '\n\n'.join(text_parts)
        
        except Exception as e:
            raise ValueError(f"Failed to extract text from PDF {filename}: {str(e)}")


class AttachmentProcessor:
    """Main processor that routes to specific handlers"""
    
    @staticmethod
    def process_attachment(file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Process attachment based on file type
        Returns dict with processed data
        """
        file_type = FileProcessor.get_file_type(filename)
        
        try:
            if file_type == 'text':
                text = TextFileProcessor.extract_text(file_content, filename)
                return {
                    'type': 'text',
                    'filename': filename,
                    'content': text,
                    'char_count': len(text)
                }
            
            elif file_type == 'image':
                # For images, prepare for GPT Vision
                return ImageProcessor.prepare_for_vision(file_content, filename)
            
            elif file_type == 'office':
                ext = Path(filename).suffix.lower()
                if ext == '.docx':
                    text = OfficeFileProcessor.extract_from_docx(file_content, filename)
                elif ext == '.xlsx':
                    text = OfficeFileProcessor.extract_from_xlsx(file_content, filename)
                elif ext == '.pptx':
                    text = OfficeFileProcessor.extract_from_pptx(file_content, filename)
                else:
                    raise ValueError(f"Unsupported office format: {ext}")
                
                return {
                    'type': 'office',
                    'filename': filename,
                    'content': text,
                    'char_count': len(text),
                    'format': ext
                }
            
            elif file_type == 'pdf':
                text = PDFProcessor.extract_text(file_content, filename)
                return {
                    'type': 'pdf',
                    'filename': filename,
                    'content': text,
                    'char_count': len(text)
                }
            
            else:
                # Unknown file type - return basic info
                return {
                    'type': 'unknown',
                    'filename': filename,
                    'size_bytes': len(file_content),
                    'content': f"[Unsupported file type: {Path(filename).suffix}]"
                }
        
        except Exception as e:
            return {
                'type': 'error',
                'filename': filename,
                'error': str(e),
                'content': f"[Error processing file: {str(e)}]"
            }
    
    @staticmethod
    def process_multiple_attachments(attachments: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Process multiple attachments
        attachments: List of dicts with 'content' (base64) and 'filename'
        Returns dict with processed attachments grouped by type
        """
        processed = {
            'text_files': [],
            'images': [],
            'office_files': [],
            'pdfs': [],
            'errors': []
        }
        
        for att in attachments:
            try:
                # Decode base64 content
                file_content = base64.b64decode(att['content'])
                filename = att['filename']
                
                # Process the file
                result = AttachmentProcessor.process_attachment(file_content, filename)
                
                # Categorize result
                if result['type'] == 'text':
                    processed['text_files'].append(result)
                elif result['type'] == 'image':
                    processed['images'].append(result)
                elif result['type'] in ['office', 'pdf']:
                    processed['office_files'].append(result)
                elif result['type'] == 'error':
                    processed['errors'].append(result)
            
            except Exception as e:
                processed['errors'].append({
                    'filename': att.get('filename', 'unknown'),
                    'error': str(e)
                })
        
        return processed
